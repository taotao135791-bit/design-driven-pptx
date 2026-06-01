#!/usr/bin/env python3
"""
PPTD to PPTX Converter
Converts .pptd presentations (YAML-based intermediate format) to .pptx files.
Requires: python-pptx, PyYAML
Optional: resvg (for SVG decoration support)

Usage:
    python pptd2pptx.py <path-to.pptd> [<output.pptx>]
"""

import argparse
import sys
import os
import re
import copy
import tempfile
import math
import urllib.request
from html.parser import HTMLParser
from typing import Optional, List, Dict, Any, Tuple

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree

# Optional resvg for SVG rendering (deprecated — prefer programmatic shapes)
try:
    import resvg
    RESVG_AVAILABLE = True
except ImportError:
    RESVG_AVAILABLE = False

# ---------------------------------------------------------------------------
# Coordinate mapping
# ---------------------------------------------------------------------------
PPTD_WIDTH = 1280.0
PPTD_HEIGHT = 720.0


def px_to_emu(val, slide_dim_emu, total_px):
    """Map PPTD px to EMU proportionally."""
    return int(val / total_px * slide_dim_emu)


# ---------------------------------------------------------------------------
# Color / Theme helpers
# ---------------------------------------------------------------------------
def hex_to_rgb(hex_color):
    """Convert #RRGGBB or #RRGGBBAA to RGBColor (ignores alpha)."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) >= 6:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)


# ---------------------------------------------------------------------------
# Color manipulation utilities (for dynamic color adaptation)
# ---------------------------------------------------------------------------
def hex_to_hsl(hex_color):
    """Convert hex to HSL tuple (h, s, l) where h in [0,360], s,l in [0,1]."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16) / 255.0
    g = int(hex_color[2:4], 16) / 255.0
    b = int(hex_color[4:6], 16) / 255.0
    mx = max(r, g, b)
    mn = min(r, g, b)
    l = (mx + mn) / 2.0
    if mx == mn:
        h = s = 0.0
    else:
        d = mx - mn
        s = d / (2.0 - mx - mn) if l > 0.5 else d / (mx + mn)
        if mx == r:
            h = (60 * ((g - b) / d) + 360) % 360
        elif mx == g:
            h = (60 * ((b - r) / d) + 120) % 360
        else:
            h = (60 * ((r - g) / d) + 240) % 360
    return (h, s, l)


def hsl_to_hex(h, s, l):
    """Convert HSL to hex color string."""
    def hue_to_rgb(p, q, t):
        if t < 0:
            t += 1
        if t > 1:
            t -= 1
        if t < 1/6:
            return p + (q - p) * 6 * t
        if t < 1/2:
            return q
        if t < 2/3:
            return p + (q - p) * (2/3 - t) * 6
        return p

    if s == 0:
        r = g = b = l
    else:
        q = l * (1 + s) if l < 0.5 else l + s - l * s
        p = 2 * l - q
        r = hue_to_rgb(p, q, h / 360 + 1/3)
        g = hue_to_rgb(p, q, h / 360)
        b = hue_to_rgb(p, q, h / 360 - 1/3)

    return "#{:02x}{:02x}{:02x}".format(
        int(round(r * 255)), int(round(g * 255)), int(round(b * 255))
    )


def adjust_saturation(hex_color, factor):
    """Adjust color saturation by factor (0.0-1.0 reduces, >1.0 increases)."""
    h, s, l = hex_to_hsl(hex_color)
    s = max(0.0, min(1.0, s * factor))
    return hsl_to_hex(h, s, l)


def adjust_lightness(hex_color, factor):
    """Adjust color lightness by factor (0.0-1.0 darkens, >1.0 lightens)."""
    h, s, l = hex_to_hsl(hex_color)
    l = max(0.0, min(1.0, l * factor))
    return hsl_to_hex(h, s, l)


def adjust_contrast(hex_color, factor):
    """Adjust color contrast: factor > 1 pushes away from mid-gray, < 1 pulls toward."""
    h, s, l = hex_to_hsl(hex_color)
    if l > 0.5:
        l = 0.5 + (l - 0.5) * factor
    else:
        l = 0.5 - (0.5 - l) * factor
    l = max(0.0, min(1.0, l))
    return hsl_to_hex(h, s, l)


class ThemeResolver:
    def __init__(self, theme_yaml):
        self.colors = theme_yaml.get('colors', {})
        self.text_styles = theme_yaml.get('textStyles', {})
        self.table_styles = theme_yaml.get('tableStyles', {})

    def resolve_color(self, color_str):
        if not color_str:
            return None
        if color_str.startswith('$'):
            key = color_str[1:]
            color_str = self.colors.get(key, color_str)
        if color_str.startswith('#'):
            return hex_to_rgb(color_str)
        return None

    def resolve_text_style(self, style_ref):
        if not style_ref:
            return {}
        if style_ref.startswith('$'):
            key = style_ref[1:]
            return copy.deepcopy(self.text_styles.get(key, {}))
        return {}

    def get_color_palette(self):
        """Extract a usable color palette from theme colors for charts.
        Returns a list of RGBColor objects."""
        palette = []
        # Priority order for chart colors
        color_keys = ['primary', 'accent', 'secondary', 'primaryDark', 'ink', 'surfaceAlt']
        for key in color_keys:
            val = self.colors.get(key)
            if val:
                color = self.resolve_color(val)
                if color and color not in palette:
                    palette.append(color)
        # If still too few, add some derived colors from primary
        if len(palette) < 3 and palette:
            primary = palette[0]
            # Add lighter/darker variants
            palette.append(RGBColor(
                min(255, int(primary[0] * 1.3)),
                min(255, int(primary[1] * 1.3)),
                min(255, int(primary[2] * 1.3))
            ))
            palette.append(RGBColor(
                max(0, int(primary[0] * 0.7)),
                max(0, int(primary[1] * 0.7)),
                max(0, int(primary[2] * 0.7))
            ))
        return palette if palette else [RGBColor(26, 26, 26)]

    def get_chart_font(self):
        """Get the best font family for chart text from theme."""
        # Prefer body style, then any style with a fontFamily
        for key in ['body', 'bodySm', 'heroTitle', 'sectionHeadline']:
            style = self.text_styles.get(key, {})
            if style.get('fontFamily'):
                return style['fontFamily'].split(',')[0].strip()
        return "Arial"


# ---------------------------------------------------------------------------
# SVG helper (legacy — prefer programmatic shapes)
# ---------------------------------------------------------------------------
def svg_to_png_bytes(svg_path):
    """Render SVG file to PNG bytes using resvg. Returns None if unavailable."""
    if not RESVG_AVAILABLE:
        return None
    try:
        with open(svg_path, 'rb') as f:
            svg_data = f.read()
        opts = resvg.usvg.Options.default()
        tree = resvg.usvg.Tree.from_str(svg_data.decode('utf-8'), opts)
        size = tree.int_size()
        transform = (2.0, 0.0, 0.0, 2.0, 0.0, 0.0)
        png = resvg.render(tree, transform)
        return png
    except (OSError, ImportError, ValueError, RuntimeError) as e:
        print(f"Warning: SVG rendering failed for {svg_path}: {e}")
        return None


# ---------------------------------------------------------------------------
# HTML Rich Text Parser
# ---------------------------------------------------------------------------
class RichTextState:
    def __init__(self):
        self.bold = False
        self.italic = False
        self.underline = False
        self.strike = False
        self.sup = False
        self.sub = False
        self.color = None
        self.font_size = None
        self.font_family = None
        self.href = None


class RichTextSegment:
    def __init__(self, text, state, para_style=None):
        self.text = text
        self.state = state
        self.para_style = para_style or {}


class RichTextParser(HTMLParser):
    """Parse PPTD rich text (subset of HTML) into segments."""

    def __init__(self):
        super().__init__()
        self.segments = []
        self.state_stack = [RichTextState()]
        self.para_style = {}
        self.current_text = []

    def _current_state(self):
        base = RichTextState()
        for st in self.state_stack:
            base.bold = base.bold or st.bold
            base.italic = base.italic or st.italic
            base.underline = base.underline or st.underline
            base.strike = base.strike or st.strike
            base.sup = base.sup or st.sup
            base.sub = base.sub or st.sub
            if st.color:
                base.color = st.color
            if st.font_size:
                base.font_size = st.font_size
            if st.font_family:
                base.font_family = st.font_family
            if st.href:
                base.href = st.href
        return base

    def _flush(self):
        if self.current_text:
            text = ''.join(self.current_text)
            if text:
                self.segments.append(RichTextSegment(
                    text, self._current_state(), dict(self.para_style)
                ))
            self.current_text = []

    def handle_starttag(self, tag, attrs):
        self._flush()
        attrs_dict = dict(attrs)
        if tag == 'p':
            self.para_style = {}
            style = attrs_dict.get('style', '')
            if 'text-align:' in style:
                m = re.search(r'text-align:\s*(left|center|right|justify)', style)
                if m:
                    self.para_style['align'] = m.group(1)
            if 'line-height:' in style:
                m = re.search(r'line-height:\s*([0-9.]+)', style)
                if m:
                    self.para_style['line_height'] = float(m.group(1))
            if 'margin-top:' in style:
                m = re.search(r'margin-top:\s*([0-9]+)px', style)
                if m:
                    self.para_style['space_before'] = int(m.group(1))
        elif tag == 'span':
            st = RichTextState()
            style = attrs_dict.get('style', '')
            if 'color:' in style:
                m = re.search(r'color:\s*([^;\s]+)', style)
                if m:
                    st.color = m.group(1)
            if 'font-size:' in style:
                m = re.search(r'font-size:\s*([0-9]+)px', style)
                if m:
                    st.font_size = int(m.group(1))
            if 'font-family:' in style:
                m = re.search(r"font-family:([^;]+)", style)
                if m:
                    st.font_family = m.group(1).strip().strip('"')
            self.state_stack.append(st)
        elif tag == 'strong':
            st = RichTextState(); st.bold = True
            self.state_stack.append(st)
        elif tag == 'em':
            st = RichTextState(); st.italic = True
            self.state_stack.append(st)
        elif tag == 'u':
            st = RichTextState(); st.underline = True
            self.state_stack.append(st)
        elif tag == 's':
            st = RichTextState(); st.strike = True
            self.state_stack.append(st)
        elif tag == 'sup':
            st = RichTextState(); st.sup = True
            self.state_stack.append(st)
        elif tag == 'sub':
            st = RichTextState(); st.sub = True
            self.state_stack.append(st)
        elif tag == 'a':
            st = RichTextState()
            st.href = attrs_dict.get('href', '')
            self.state_stack.append(st)
        elif tag in ('ul', 'ol'):
            self.para_style['bullet'] = True
            self.para_style['bullet_type'] = 'bullet' if tag == 'ul' else 'number'
        elif tag == 'br':
            self.current_text.append('\n')

    def handle_endtag(self, tag):
        self._flush()
        if tag in ('p', 'span', 'strong', 'em', 'u', 's', 'sup', 'sub', 'a'):
            if len(self.state_stack) > 1:
                self.state_stack.pop()
        elif tag in ('ul', 'ol'):
            self.para_style['bullet'] = False

    def handle_data(self, data):
        self.current_text.append(data)

    def get_segments(self):
        self._flush()
        return self.segments


def parse_rich_text(html_text):
    """Parse rich text string into list of RichTextSegment."""
    if '<' not in html_text:
        html_text = f'<p>{html_text}</p>'
    parser = RichTextParser()
    parser.feed(html_text)
    return parser.get_segments()


# ---------------------------------------------------------------------------
# Alignment mapping
# ---------------------------------------------------------------------------
H_ALIGN_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY,
    'distributed': PP_ALIGN.DISTRIBUTE,
}

V_ALIGN_MAP = {
    'top': MSO_ANCHOR.TOP,
    'middle': MSO_ANCHOR.MIDDLE,
    'bottom': MSO_ANCHOR.BOTTOM,
}


# ---------------------------------------------------------------------------
# Shape name mapping
# ---------------------------------------------------------------------------
SHAPE_MAP = {
    'rect': MSO_SHAPE.RECTANGLE,
    'roundRect': MSO_SHAPE.ROUNDED_RECTANGLE,
    'ellipse': MSO_SHAPE.OVAL,
    'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    'diamond': MSO_SHAPE.DIAMOND,
    'plus': MSO_SHAPE.CROSS,
    'star5': MSO_SHAPE.STAR_5_POINT,
    'rightArrow': MSO_SHAPE.RIGHT_ARROW,
    'homePlate': MSO_SHAPE.PENTAGON,
    'chevron': MSO_SHAPE.CHEVRON,
    'donut': MSO_SHAPE.DONUT,
    'flowChartProcess': MSO_SHAPE.FLOWCHART_PROCESS,
    'flowChartDecision': MSO_SHAPE.FLOWCHART_DECISION,
    'flowChartTerminator': MSO_SHAPE.FLOWCHART_TERMINATOR,
}


# ---------------------------------------------------------------------------
# Chart type mapping (PPTD -> python-pptx)
# ---------------------------------------------------------------------------
CHART_TYPE_MAP = {
    'bar': {'vertical': 'COLUMN_CLUSTERED', 'horizontal': 'BAR_CLUSTERED'},
    'line': 'LINE_MARKERS',
    'area': 'AREA',
    'scatter': 'XY_SCATTER',
    'pie': 'PIE',
    'radar': 'RADAR',
    'combo': 'COLUMN_CLUSTERED',
    'bubble': 'BUBBLE',
}


# ---------------------------------------------------------------------------
# Converter
# ---------------------------------------------------------------------------
class PPTDConverter:
    def __init__(self, pptd_path):
        self.pptd_path = pptd_path
        self.base_dir = os.path.dirname(os.path.abspath(pptd_path))
        with open(pptd_path, encoding='utf-8') as f:
            self.pptd = yaml.safe_load(f)
        self.theme = ThemeResolver(self.pptd.get('theme', {}))
        size = self.pptd.get('size', [1280, 720])
        self.pptd_w, self.pptd_h = size[0], size[1]
        self.prs = Presentation()
        self.slide_width_emu = 12192000
        self.slide_height_emu = 6858000
        self.prs.slide_width = Emu(self.slide_width_emu)
        self.prs.slide_height = Emu(self.slide_height_emu)

    def _to_emu(self, x, y, w, h):
        """Convert PPTD bounds [x,y,w,h] to EMU."""
        x_emu = px_to_emu(x, self.slide_width_emu, self.pptd_w)
        y_emu = px_to_emu(y, self.slide_height_emu, self.pptd_h)
        w_emu = px_to_emu(w, self.slide_width_emu, self.pptd_w)
        h_emu = px_to_emu(h, self.slide_height_emu, self.pptd_h)
        return x_emu, y_emu, w_emu, h_emu

    # ---------------------------------------------------------------------------
    # Color contrast helpers
    # ---------------------------------------------------------------------------
    @staticmethod
    def _color_luminance(rgb_color):
        """Calculate relative luminance of an RGBColor."""
        if rgb_color is None:
            return 0.0
        def _c(v):
            v = v / 255.0
            return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
        return 0.2126 * _c(rgb_color[0]) + 0.7152 * _c(rgb_color[1]) + 0.0722 * _c(rgb_color[2])

    def _ensure_contrast(self, rgb_color, bg_color, target=3.0):
        """Adjust a color (lighten or darken) until it achieves target contrast against bg_color.
        Returns adjusted RGBColor."""
        if rgb_color is None or bg_color is None:
            return rgb_color
        fg_lum = self._color_luminance(rgb_color)
        bg_lum = self._color_luminance(bg_color)
        ratio = (max(fg_lum, bg_lum) + 0.05) / (min(fg_lum, bg_lum) + 0.05)
        if ratio >= target:
            return rgb_color
        r, g, b = rgb_color[0], rgb_color[1], rgb_color[2]
        # Determine direction: if background is bright, darken the color; if dark, lighten it
        if bg_lum > 0.5:
            # Darken
            for _ in range(20):
                r = max(0, int(r * 0.85) - 5)
                g = max(0, int(g * 0.85) - 5)
                b = max(0, int(b * 0.85) - 5)
                test = RGBColor(r, g, b)
                test_lum = self._color_luminance(test)
                new_ratio = (max(test_lum, bg_lum) + 0.05) / (min(test_lum, bg_lum) + 0.05)
                if new_ratio >= target:
                    return test
        else:
            # Lighten
            for _ in range(20):
                r = min(255, int(r * 1.15) + 5)
                g = min(255, int(g * 1.15) + 5)
                b = min(255, int(b * 1.15) + 5)
                test = RGBColor(r, g, b)
                test_lum = self._color_luminance(test)
                new_ratio = (max(test_lum, bg_lum) + 0.05) / (min(test_lum, bg_lum) + 0.05)
                if new_ratio >= target:
                    return test
        return RGBColor(r, g, b)

    # =====================================================================
    # Animation support
    # =====================================================================
    class _AnimationBuilder:
        """Builds PPTX animation XML via lxml."""

        # Animation type → (filter_prefix, default_transition)
        ANIM_FILTERS = {
            'fade': ('fade', 'in'),
            'appear': ('appear', 'in'),
            'disappear': ('disappear', 'out'),
            'wipe': ('wipe', 'in'),
            'peekIn': ('peekIn', 'in'),
            'peekOut': ('peekOut', 'out'),
            'flyIn': ('flyIn', 'in'),
            'flyOut': ('flyOut', 'out'),
            'growShrink': ('growShrink', 'in'),
            'spin': ('spin', 'in'),
            'bounce': ('bounce', 'in'),
            'pulse': ('pulse', 'in'),
            'swivel': ('swivel', 'in'),
            'float': ('float', 'in'),
            'split': ('split', 'in'),
            'stripe': ('stripe', 'in'),
            'compress': ('compress', 'in'),
        }

        # Direction normalization for filter params
        DIR_MAP = {'left': 'left', 'right': 'right', 'top': 'up', 'bottom': 'down',
                   'up': 'up', 'down': 'down', 'center': 'center', 'across': 'across'}

        TRIGGER_NODE_TYPES = {
            'onClick': 'clickEffect',
            'withPrevious': 'withEffect',
            'afterPrevious': 'afterEffect',
        }

        def __init__(self, slide):
            self.slide = slide
            self._next_id = 10
            self._seq_childTnLst = None

        def _get_id(self):
            self._next_id += 1
            return str(self._next_id)

        def _ensure_timing(self):
            """Create <p:timing> structure on the slide if absent."""
            if self._seq_childTnLst is not None:
                return
            slide_el = self.slide._element
            NS = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            timing = slide_el.find(qn('p:timing'))
            if timing is None:
                timing = etree.SubElement(slide_el, qn('p:timing'))
                tnLst = etree.SubElement(timing, qn('p:tnLst'))
                outer_par = etree.SubElement(tnLst, qn('p:par'))
                outer_ctn = etree.SubElement(outer_par, qn('p:cTn'))
                outer_ctn.set('id', self._get_id())
                outer_ctn.set('dur', 'indefinite')
                outer_ctn.set('restart', 'never')
                outer_ctn.set('nodeType', 'tmRoot')
                childTnLst = etree.SubElement(outer_ctn, qn('p:childTnLst'))
                seq = etree.SubElement(childTnLst, qn('p:seq'))
                seq.set('concurrent', '1')
                seq.set('nextAc', 'seek')
                seq_ctn = etree.SubElement(seq, qn('p:cTn'))
                seq_ctn.set('id', self._get_id())
                seq_ctn.set('dur', 'indefinite')
                seq_ctn.set('nodeType', 'mainSeq')
                self._seq_childTnLst = etree.SubElement(seq_ctn, qn('p:childTnLst'))
            else:
                # Try to locate existing mainSeq childTnLst
                seq = timing.find('.//p:seq', namespaces=NS)
                if seq is not None:
                    main_ctn = seq.find('.//p:cTn[@nodeType="mainSeq"]', namespaces=NS)
                    if main_ctn is not None:
                        self._seq_childTnLst = main_ctn.find(qn('p:childTnLst'))
                if self._seq_childTnLst is None:
                    # Rebuild from scratch
                    timing.clear()
                    tnLst = etree.SubElement(timing, qn('p:tnLst'))
                    outer_par = etree.SubElement(tnLst, qn('p:par'))
                    outer_ctn = etree.SubElement(outer_par, qn('p:cTn'))
                    outer_ctn.set('id', self._get_id())
                    outer_ctn.set('dur', 'indefinite')
                    outer_ctn.set('restart', 'never')
                    outer_ctn.set('nodeType', 'tmRoot')
                    childTnLst = etree.SubElement(outer_ctn, qn('p:childTnLst'))
                    seq = etree.SubElement(childTnLst, qn('p:seq'))
                    seq.set('concurrent', '1')
                    seq.set('nextAc', 'seek')
                    seq_ctn = etree.SubElement(seq, qn('p:cTn'))
                    seq_ctn.set('id', self._get_id())
                    seq_ctn.set('dur', 'indefinite')
                    seq_ctn.set('nodeType', 'mainSeq')
                    self._seq_childTnLst = etree.SubElement(seq_ctn, qn('p:childTnLst'))

        def _build_filter(self, anim_type, direction):
            prefix_info = self.ANIM_FILTERS.get(anim_type)
            if not prefix_info:
                return None
            prefix, _ = prefix_info
            dir_str = self.DIR_MAP.get(direction, '')
            if dir_str and prefix in ('wipe', 'peekIn', 'peekOut', 'flyIn', 'flyOut'):
                return f'{prefix}({dir_str})'
            return prefix

        def add_animation(self, shape, anim_spec):
            """Add an animation effect targeting *shape*.

            anim_spec fields:
                type:     fade | appear | disappear | wipe | peekIn | peekOut | flyIn | flyOut | growShrink
                direction: left | right | top | bottom | up | down
                trigger:  onClick | withPrevious | afterPrevious
                duration: milliseconds (default 500)
                delay:    milliseconds (default 0)
            """
            self._ensure_timing()
            anim_type = anim_spec.get('type', 'fade')
            direction = anim_spec.get('direction', '')
            trigger = anim_spec.get('trigger', 'onClick')
            duration = anim_spec.get('duration', 500)
            delay = anim_spec.get('delay', 0)

            filter_str = self._build_filter(anim_type, direction)
            if not filter_str:
                return

            # Determine transition direction
            prefix_info = self.ANIM_FILTERS.get(anim_type, ('fade', 'in'))
            _, default_trans = prefix_info
            transition = 'out' if 'Out' in anim_type or 'disappear' in anim_type else default_trans

            # Trigger / node type
            node_type = self.TRIGGER_NODE_TYPES.get(trigger, 'clickEffect')
            delay_str = 'indefinite' if trigger == 'onClick' else str(delay)

            # Build XML
            par = etree.SubElement(self._seq_childTnLst, qn('p:par'))
            ctn = etree.SubElement(par, qn('p:cTn'))
            ctn.set('id', self._get_id())
            ctn.set('fill', 'hold')
            ctn.set('nodeType', node_type)
            stCondLst = etree.SubElement(ctn, qn('p:stCondLst'))
            cond = etree.SubElement(stCondLst, qn('p:cond'))
            cond.set('delay', delay_str)
            ctn_childTnLst = etree.SubElement(ctn, qn('p:childTnLst'))

            anim = etree.SubElement(ctn_childTnLst, qn('p:animEffect'))
            anim.set('transition', transition)
            anim.set('filter', filter_str)
            anim.set('dur', str(duration))
            cBhvr = etree.SubElement(anim, qn('p:cBhvr'))
            cBhvr_ctn = etree.SubElement(cBhvr, qn('p:cTn'))
            cBhvr_ctn.set('id', self._get_id())
            cBhvr_ctn.set('dur', str(duration))
            cBhvr_ctn.set('fill', 'hold')
            tgtEl = etree.SubElement(cBhvr, qn('p:tgtEl'))
            spTgt = etree.SubElement(tgtEl, qn('p:spTgt'))
            spTgt.set('spid', str(shape.shape_id))

    def _apply_animation(self, slide, shape, element):
        """Apply animation from element spec to shape if present."""
        anim_spec = element.get('animation')
        if not anim_spec:
            return
        builder = getattr(slide, '_pptd_anim_builder', None)
        if builder is None:
            builder = self._AnimationBuilder(slide)
            slide._pptd_anim_builder = builder
        builder.add_animation(shape, anim_spec)

    def _resolve_style(self, content):
        """Merge text style: theme style + inline overrides."""
        style = {}
        style_ref = content.get('style')
        if style_ref:
            style = self.theme.resolve_text_style(style_ref)
        for key in ('fontSize', 'color', 'fontFamily', 'lineHeight',
                    'lineHeightPx', 'letterSpacing', 'marginTop'):
            if key in content:
                style[key] = content[key]
        return style

    # =====================================================================
    # NEW: Freeform / Path Support
    # =====================================================================
    def _parse_path_commands(self, path_cmds):
        """Parse SVG-like path string into list of (cmd, [(x,y), ...])."""
        tokens = re.findall(r'([MLCQZHVThvmlcqz])([^MLCQZHVThvmlcqz]*)', path_cmds)
        result = []
        for cmd, coords_str in tokens:
            # Normalize: replace commas with spaces, then split
            normalized = coords_str.replace(',', ' ')
            coords = [float(c) for c in normalized.strip().split() if c.strip()]
            pts = []
            if cmd in ('M', 'L'):
                for i in range(0, len(coords), 2):
                    if i + 1 < len(coords):
                        pts.append((coords[i], coords[i + 1]))
            elif cmd == 'C':
                for i in range(0, len(coords), 6):
                    if i + 5 < len(coords):
                        pts.append((coords[i], coords[i+1], coords[i+2], coords[i+3], coords[i+4], coords[i+5]))
            elif cmd == 'Q':
                for i in range(0, len(coords), 4):
                    if i + 3 < len(coords):
                        pts.append((coords[i], coords[i+1], coords[i+2], coords[i+3]))
            elif cmd in ('H', 'V'):
                for c in coords:
                    pts.append((c,))
            elif cmd in ('Z', 'z'):
                pts = []
            result.append((cmd, pts))
        return result

    def _sample_cubic_bezier(self, p0, p1, p2, p3, steps=12):
        """Sample points along a cubic Bezier curve."""
        pts = []
        for i in range(1, steps + 1):
            t = i / steps
            t2 = t * t
            t3 = t2 * t
            mt = 1 - t
            mt2 = mt * mt
            mt3 = mt2 * mt
            x = mt3 * p0[0] + 3 * mt2 * t * p1[0] + 3 * mt * t2 * p2[0] + t3 * p3[0]
            y = mt3 * p0[1] + 3 * mt2 * t * p1[1] + 3 * mt * t2 * p2[1] + t3 * p3[1]
            pts.append((int(round(x)), int(round(y))))
        return pts

    def _sample_quadratic_bezier(self, p0, p1, p2, steps=10):
        """Sample points along a quadratic Bezier curve."""
        pts = []
        for i in range(1, steps + 1):
            t = i / steps
            mt = 1 - t
            x = mt * mt * p0[0] + 2 * mt * t * p1[0] + t * t * p2[0]
            y = mt * mt * p0[1] + 2 * mt * t * p1[1] + t * t * p2[1]
            pts.append((int(round(x)), int(round(y))))
        return pts

    def _add_freeform_shape(self, slide, element):
        """Draw a custom freeform shape from SVG-like path data.
        
        Path format: "viewW,viewH;M x y L x y ... Z"
        Coordinates are in viewBox space and mapped to element bounds.
        Cubic (C) and Quadratic (Q) Bezier curves are sampled into line segments.
        """
        path_str = element.get('path', '')
        if not path_str:
            return None

        parts = path_str.split(';')
        if len(parts) < 2:
            return None

        view_size = parts[0].strip().split(',')
        if len(view_size) != 2:
            return None

        view_w, view_h = float(view_size[0]), float(view_size[1])
        path_cmds = parts[1].strip()

        bounds = element.get('bounds', [0, 0, 100, 100])
        if len(bounds) < 4:
            return None
        bx, by, bw, bh = bounds
        bx_emu, by_emu, bw_emu, bh_emu = self._to_emu(bx, by, bw, bh)

        def map_x(x):
            return bx_emu + int(x / view_w * bw_emu)

        def map_y(y):
            return by_emu + int(y / view_h * bh_emu)

        cmds = self._parse_path_commands(path_cmds)
        if not cmds:
            return None

        # Find start point (first M)
        start_pt = None
        for cmd, pts in cmds:
            if cmd == 'M' and pts:
                start_pt = (map_x(pts[0][0]), map_y(pts[0][1]))
                break

        if start_pt is None:
            return None

        builder = slide.shapes.build_freeform(Emu(start_pt[0]), Emu(start_pt[1]))

        current_pt = (float(start_pt[0]), float(start_pt[1]))
        for cmd, pts in cmds:
            if cmd == 'M':
                for i, pt in enumerate(pts):
                    if i == 0:
                        continue  # start point already set
                    px, py = map_x(pt[0]), map_y(pt[1])
                    builder.add_line_segments([(Emu(px), Emu(py))])
                    current_pt = (float(px), float(py))
            elif cmd == 'L':
                segs = []
                for pt in pts:
                    px, py = map_x(pt[0]), map_y(pt[1])
                    segs.append((Emu(px), Emu(py)))
                    current_pt = (float(px), float(py))
                if segs:
                    builder.add_line_segments(segs)
            elif cmd == 'C':
                for pt in pts:
                    p0 = current_pt
                    p1 = (float(map_x(pt[0])), float(map_y(pt[1])))
                    p2 = (float(map_x(pt[2])), float(map_y(pt[3])))
                    p3 = (float(map_x(pt[4])), float(map_y(pt[5])))
                    sampled = self._sample_cubic_bezier(p0, p1, p2, p3, steps=12)
                    for px, py in sampled:
                        builder.add_line_segments([(Emu(px), Emu(py))])
                    current_pt = p3
            elif cmd == 'Q':
                for pt in pts:
                    p0 = current_pt
                    p1 = (float(map_x(pt[0])), float(map_y(pt[1])))
                    p2 = (float(map_x(pt[2])), float(map_y(pt[3])))
                    sampled = self._sample_quadratic_bezier(p0, p1, p2, steps=10)
                    for px, py in sampled:
                        builder.add_line_segments([(Emu(px), Emu(py))])
                    current_pt = p2
            elif cmd == 'H':
                for pt in pts:
                    px = map_x(pt[0])
                    py = int(current_pt[1])
                    builder.add_line_segments([(Emu(px), Emu(py))])
                    current_pt = (float(px), float(py))
            elif cmd == 'V':
                for pt in pts:
                    px = int(current_pt[0])
                    py = map_y(pt[0])
                    builder.add_line_segments([(Emu(px), Emu(py))])
                    current_pt = (float(px), float(py))
            elif cmd in ('Z', 'z'):
                builder._add_close()

        shape = builder.convert_to_shape()
        rotation = element.get('rotation', 0)
        if rotation:
            shape.rotation = rotation

        fill_spec = element.get('fill')
        if fill_spec:
            self._apply_fill(shape, fill_spec)
        else:
            shape.fill.background()

        border_spec = element.get('border')
        self._apply_border(shape, border_spec)

        shadow_spec = element.get('shadow')
        if shadow_spec:
            self._apply_shadow(shape, shadow_spec)

        opacity = element.get('opacity', 1)
        if opacity < 1:
            self._apply_opacity(shape, opacity)

        self._apply_animation(slide, shape, element)
        return shape

    # =====================================================================
    # NEW: Shadow Support via DrawingML
    # =====================================================================
    def _apply_shadow(self, shape, shadow_spec):
        """Apply outer shadow via DrawingML XML."""
        if not shadow_spec:
            return

        color_str = shadow_spec.get('color')
        if not color_str:
            return

        color = self.theme.resolve_color(color_str)
        if not color:
            return

        offset = shadow_spec.get('offset')
        if offset and isinstance(offset, (list, tuple)) and len(offset) >= 2:
            dx, dy = offset[0], offset[1]
        else:
            dx = shadow_spec.get('dx', shadow_spec.get('offsetX', 4))
            dy = shadow_spec.get('dy', shadow_spec.get('offsetY', 4))
        blur = shadow_spec.get('blur', 0)

        # Convert dx/dy to PowerPoint distance+direction
        dist = math.sqrt(dx * dx + dy * dy)
        angle = math.degrees(math.atan2(dy, dx))
        direction = int((angle % 360) * 60000)
        distance = int(dist * 12700)  # 1pt = 12700 EMU

        spPr = shape._element.spPr

        # Remove existing effectLst
        existing = spPr.find(qn('a:effectLst'))
        if existing is not None:
            spPr.remove(existing)

        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(int(blur * 12700)))
        outerShdw.set('dist', str(distance))
        outerShdw.set('dir', str(direction))
        outerShdw.set('sx', '100000')
        outerShdw.set('sy', '100000')

        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')

    # =====================================================================
    # NEW: Opacity Support via DrawingML alpha
    # =====================================================================
    def _apply_opacity(self, shape, opacity):
        """Set fill opacity via alpha element in DrawingML."""
        if opacity >= 1:
            return

        spPr = shape._element.spPr
        alpha_val = max(0, min(100000, int(opacity * 100000)))

        # Try solidFill
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            for color_elem in list(solidFill):
                alpha = color_elem.find(qn('a:alpha'))
                if alpha is None:
                    alpha = etree.SubElement(color_elem, qn('a:alpha'))
                alpha.set('val', str(alpha_val))
            return

        # Try gradFill stops
        gradFill = spPr.find(qn('a:gradFill'))
        if gradFill is not None:
            gsLst = gradFill.find(qn('a:gsLst'))
            if gsLst is not None:
                for gs in gsLst:
                    for color_elem in list(gs):
                        alpha = color_elem.find(qn('a:alpha'))
                        if alpha is None:
                            alpha = etree.SubElement(color_elem, qn('a:alpha'))
                        alpha.set('val', str(alpha_val))

    # =====================================================================
    # ENHANCED: Fill with true multi-stop gradient
    # =====================================================================
    def _apply_fill(self, shape, fill_spec):
        if not fill_spec:
            return
        fill_type = fill_spec.get('type')
        if fill_type == 'solid':
            color = self.theme.resolve_color(fill_spec.get('color'))
            if color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
        elif fill_type == 'gradient':
            grad_type = fill_spec.get('gradientType', 'linear')
            stops = fill_spec.get('stops', [])
            angle = fill_spec.get('angle', 0)

            if stops and len(stops) >= 2:
                shape.fill.gradient()
                fill_elem = shape._element.spPr.find(qn('a:gradFill'))
                if fill_elem is not None:
                    # Set linear gradient angle
                    if grad_type == 'linear':
                        lin = fill_elem.find(qn('a:lin'))
                        if lin is None:
                            lin = etree.SubElement(fill_elem, qn('a:lin'))
                        # PowerPoint: 0=right, 90=down, etc. In 1/60000 degrees
                        lin.set('ang', str(int((angle % 360) * 60000)))
                        lin.set('scaled', '0')

                    # Build gradient stops
                    gsLst = fill_elem.find(qn('a:gsLst'))
                    if gsLst is not None:
                        fill_elem.remove(gsLst)
                    gsLst = etree.SubElement(fill_elem, qn('a:gsLst'))

                    for stop in stops:
                        pos = stop.get('position', 0)
                        color_str = stop.get('color', '#000000')
                        color = self.theme.resolve_color(color_str)

                        gs = etree.SubElement(gsLst, qn('a:gs'))
                        gs.set('pos', str(max(0, min(100000, int(pos * 100000)))))

                        srgbClr = etree.SubElement(gs, qn('a:srgbClr'))
                        srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')

                    # Ensure tile rect doesn't break things
                    tileRect = fill_elem.find(qn('a:tileRect'))
                    if tileRect is not None:
                        fill_elem.remove(tileRect)
        elif fill_type == 'image':
            src = fill_spec.get('src', '')
            if src:
                print(f"Warning: Image fill on shapes is not fully supported. Use elementType: image instead.")

    def _apply_border(self, shape, border_spec):
        if not border_spec:
            return
        width = border_spec.get('width', 1)
        color = self.theme.resolve_color(border_spec.get('color'))
        if color:
            shape.line.color.rgb = color
            shape.line.width = Pt(width)
        else:
            shape.line.fill.background()

    def _add_text_frame(self, shape, content):
        """Add rich text to a shape's text frame."""
        tf = shape.text_frame
        tf.word_wrap = content.get('wrap', True)
        tf.auto_size = None

        align = content.get('align', ['left', 'top'])
        h_align = H_ALIGN_MAP.get(align[0], PP_ALIGN.LEFT)
        v_align = V_ALIGN_MAP.get(align[1], MSO_ANCHOR.TOP)
        tf.paragraphs[0].alignment = h_align

        base_style = self._resolve_style(content)
        html_text = content.get('text', '')
        segments = parse_rich_text(html_text)

        if not segments:
            return

        current_para = tf.paragraphs[0]
        first = True

        # Apply base lineHeight to first paragraph before segments
        base_line_height = base_style.get('lineHeight')
        if base_line_height is not None:
            tf.paragraphs[0].line_spacing = base_line_height

        for seg in segments:
            if '\n' in seg.text:
                parts = seg.text.split('\n')
                for i, part in enumerate(parts):
                    if i > 0:
                        current_para = tf.add_paragraph()
                        current_para.alignment = h_align
                        if base_line_height is not None:
                            current_para.line_spacing = base_line_height
                    if part:
                        run = current_para.add_run()
                        run.text = part
                        self._apply_run_style(run, seg.state, base_style)
            else:
                run = current_para.add_run()
                run.text = seg.text
                self._apply_run_style(run, seg.state, base_style)

            if seg.para_style.get('space_before'):
                current_para.space_before = Pt(seg.para_style['space_before'])
            if seg.para_style.get('line_height'):
                current_para.line_spacing = seg.para_style['line_height']

            first = False

        if len(tf.paragraphs) > 1 and not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
            p._element.getparent().remove(p._element)

    def _apply_run_style(self, run, state, base_style):
        """Apply font styling to a run."""
        font = run.font
        font.bold = state.bold or base_style.get('fontBold', False)
        font.italic = state.italic or base_style.get('fontItalic', False)
        if state.underline or base_style.get('fontUnderline', False):
            font.underline = True

        color_str = state.color or base_style.get('color')
        color = self.theme.resolve_color(color_str)
        if color:
            font.color.rgb = color

        size = state.font_size or base_style.get('fontSize')
        if size:
            font.size = Pt(size)

        family = state.font_family or base_style.get('fontFamily')
        if family:
            font.name = family.split(',')[0].strip()
            # CJK font fallback via a:ea element
            if ',' in family:
                cjk_font = family.split(',')[1].strip()
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = etree.SubElement(run._r, qn('a:rPr'))
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = etree.SubElement(rPr, qn('a:ea'))
                ea.set('typeface', cjk_font)

        if state.sup:
            font._element.set(qn('a:baseline'), '30000')
        elif state.sub:
            font._element.set(qn('a:baseline'), '-30000')

    def _add_shape(self, slide, element):
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)
        shape_name = element.get('shapeName', 'rect')
        rotation = element.get('rotation', 0)

        # NEW: custom path support
        if shape_name == 'custom':
            result = self._add_freeform_shape(slide, element)
            if result is not None:
                return result
            # fallback to rectangle
            shape_name = 'rect'

        mso_shape = SHAPE_MAP.get(shape_name, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso_shape, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.rotation = rotation

        fill_spec = element.get('fill')
        if fill_spec:
            self._apply_fill(shape, fill_spec)
        else:
            shape.fill.background()

        border_spec = element.get('border')
        self._apply_border(shape, border_spec)

        shadow_spec = element.get('shadow')
        if shadow_spec:
            self._apply_shadow(shape, shadow_spec)

        # NEW: opacity support
        opacity = element.get('opacity', 1)
        if opacity < 1:
            self._apply_opacity(shape, opacity)

        self._apply_animation(slide, shape, element)
        return shape

    # =====================================================================
    # NEW: Programmatic decoration helpers
    # =====================================================================
    def _draw_corner_brackets(self, slide, bounds, color, stroke_width=4, size=24, style='L'):
        """Draw L-shaped corner brackets at opposite corners.
        
        style='L': top-left and bottom-right (outward)
        style='frame': all four corners
        """
        bx, by, bw, bh = bounds
        sw = stroke_width

        def make_bracket(x, y, flip_x=False, flip_y=False):
            # L shape: horizontal arm + vertical arm
            # ViewBox 24x24, path: M0 0 L24 0 L24 4 L4 4 L4 24 L0 24 Z
            # This creates an L shape with thickness 4
            if not flip_x and not flip_y:
                # Top-left outward
                path = f"{size},{size};M0 0 L{size} 0 L{size} {sw} L{sw} {sw} L{sw} {size} L0 {size} Z"
                return [x, y, size, size, path]
            elif flip_x and not flip_y:
                # Top-right outward
                path = f"{size},{size};M{size} 0 L0 0 L0 {sw} L{size-sw} {sw} L{size-sw} {size} L{size} {size} Z"
                return [x + bw - size, y, size, size, path]
            elif not flip_x and flip_y:
                # Bottom-left outward
                path = f"{size},{size};M0 {size} L{size} {size} L{size} {size-sw} L{sw} {size-sw} L{sw} 0 L0 0 Z"
                return [x, y + bh - size, size, size, path]
            else:
                # Bottom-right outward
                path = f"{size},{size};M{size} {size} L0 {size} L0 {size-sw} L{size-sw} {size-sw} L{size-sw} 0 L{size} 0 Z"
                return [x + bw - size, y + bh - size, size, size, path]

        brackets = []
        if style == 'L':
            brackets = [
                make_bracket(bx, by, False, False),
                make_bracket(bx, by, True, True),
            ]
        elif style == 'frame':
            brackets = [
                make_bracket(bx, by, False, False),
                make_bracket(bx, by, True, False),
                make_bracket(bx, by, False, True),
                make_bracket(bx, by, True, True),
            ]

        for bx2, by2, bw2, bh2, path in brackets:
            elem = {
                'elementType': 'shape',
                'shapeName': 'custom',
                'path': path,
                'bounds': [bx2, by2, bw2, bh2],
                'fill': {'type': 'solid', 'color': color},
            }
            self._add_shape(slide, elem)

    def _draw_accent_line(self, slide, bounds, color, width=80, height=4, shadow_color=None, shadow_dx=4, shadow_dy=4):
        """Draw a solid accent line rectangle."""
        x, y, _, _ = bounds
        elem = {
            'elementType': 'shape',
            'shapeName': 'rect',
            'bounds': [x, y, width, height],
            'fill': {'type': 'solid', 'color': color},
        }
        if shadow_color:
            elem['shadow'] = {'color': shadow_color, 'dx': shadow_dx, 'dy': shadow_dy, 'blur': 0}
        self._add_shape(slide, elem)

    def _draw_hatch_pattern(self, slide, bounds, color, angle=45, spacing=20, line_width=2, opacity=0.06):
        """Draw diagonal hatch pattern using multiple thin rectangles."""
        bx, by, bw, bh = bounds
        import math
        rad = math.radians(angle)
        # Diagonal line length to cover the rect
        diag = math.sqrt(bw * bw + bh * bh)
        count = int(diag / spacing) + 5

        cx = bx + bw / 2
        cy = by + bh / 2

        for i in range(-count, count + 1):
            # Center point along perpendicular direction
            offset = i * spacing
            # Line center in local coords
            lx = offset * math.cos(rad + math.pi / 2)
            ly = offset * math.sin(rad + math.pi / 2)
            # Global coords
            gx = cx + lx
            gy = cy + ly
            # Line endpoints along angle direction
            dx = diag / 2 * math.cos(rad)
            dy = diag / 2 * math.sin(rad)
            x1, y1 = gx - dx, gy - dy
            x2, y2 = gx + dx, gy + dy
            # Bounding rect for this line segment
            min_x = max(bx, min(x1, x2) - line_width / 2)
            min_y = max(by, min(y1, y2) - line_width / 2)
            max_x = min(bx + bw, max(x1, x2) + line_width / 2)
            max_y = min(by + bh, max(y1, y2) + line_width / 2)
            if max_x <= min_x or max_y <= min_y:
                continue
            # Draw as thin rectangle rotated to angle
            elem = {
                'elementType': 'shape',
                'shapeName': 'rect',
                'bounds': [min_x, min_y, max_x - min_x, max_y - min_y],
                'fill': {'type': 'solid', 'color': color},
                'opacity': opacity,
                'rotation': angle,
            }
            self._add_shape(slide, elem)

    def _draw_scanlines(self, slide, bounds, color, spacing=4, line_height=2, opacity=0.04):
        """Draw horizontal CRT scanline overlay."""
        bx, by, bw, bh = bounds
        y = by
        while y < by + bh:
            elem = {
                'elementType': 'shape',
                'shapeName': 'rect',
                'bounds': [bx, y, bw, line_height],
                'fill': {'type': 'solid', 'color': color},
                'opacity': opacity,
            }
            self._add_shape(slide, elem)
            y += spacing

    def _draw_grid(self, slide, bounds, color, grid_size=40, line_width=1, opacity=0.07):
        """Draw a square grid pattern."""
        bx, by, bw, bh = bounds
        # Vertical lines
        x = bx
        while x <= bx + bw:
            elem = {
                'elementType': 'shape',
                'shapeName': 'rect',
                'bounds': [x, by, line_width, bh],
                'fill': {'type': 'solid', 'color': color},
                'opacity': opacity,
            }
            self._add_shape(slide, elem)
            x += grid_size
        # Horizontal lines
        y = by
        while y <= by + bh:
            elem = {
                'elementType': 'shape',
                'shapeName': 'rect',
                'bounds': [bx, y, bw, line_width],
                'fill': {'type': 'solid', 'color': color},
                'opacity': opacity,
            }
            self._add_shape(slide, elem)
            y += grid_size

    def _draw_decorative_number(self, slide, bounds, number, color, font_family, font_size, opacity=0.12):
        """Draw an oversized decorative numeral at low opacity."""
        x, y, w, h = bounds
        elem = {
            'elementType': 'text',
            'bounds': [x, y, w, h],
            'content': {
                'text': f'<p>{number}</p>',
                'fontSize': font_size,
                'color': color,
                'fontFamily': font_family,
                'align': ['left', 'top'],
                'wrap': False,
            },
            'opacity': opacity,
        }
        # Note: text opacity not directly supported; we'll create a group-like effect
        # by drawing a text box. For now, this is best-effort.
        self._add_text(slide, elem)

    # =====================================================================
    # Element handlers
    # =====================================================================
    def _add_text(self, slide, element):
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)

        shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        content = element.get('content', {})
        self._add_text_frame(shape, content)

        # NEW: opacity for text boxes (via run color alpha)
        opacity = element.get('opacity', 1)
        if opacity < 1:
            alpha_val = max(0, min(100000, int(opacity * 100000)))
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is None:
                        rPr = etree.SubElement(run._r, qn('a:rPr'))
                    solidFill = rPr.find(qn('a:solidFill'))
                    if solidFill is None:
                        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is None:
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', '000000')
                    alpha = srgbClr.find(qn('a:alpha'))
                    if alpha is None:
                        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                    alpha.set('val', str(alpha_val))

        self._apply_animation(slide, shape, element)
        return shape

    def _add_image(self, slide, element):
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)
        src = element.get('src', '')

        if not src:
            return None

        # Handle network URLs
        tmp_path = None
        if src.startswith('http://') or src.startswith('https://'):
            try:
                fd, tmp_path = tempfile.mkstemp(suffix=os.path.splitext(src.split('?')[0])[1] or '.jpg')
                with os.fdopen(fd, 'wb') as f:
                    req = urllib.request.Request(src, headers={'User-Agent': 'Mozilla/5.0'})
                    with urllib.request.urlopen(req, timeout=15) as resp:
                        f.write(resp.read())
                src = tmp_path
            except (urllib.error.URLError, urllib.error.HTTPError, OSError, TimeoutError) as e:
                print(f"Warning: failed to download image from {src}: {e}")
                if tmp_path and os.path.exists(tmp_path):
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
                return None
        elif not os.path.isabs(src):
            src = os.path.join(self.base_dir, src)

        if not os.path.exists(src):
            print(f"Warning: image not found: {src}")
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)
            return None

        try:
            if src.lower().endswith('.svg'):
                png_bytes = svg_to_png_bytes(src)
                if png_bytes:
                    fd, svg_tmp = tempfile.mkstemp(suffix='.png')
                    try:
                        with os.fdopen(fd, 'wb') as f:
                            f.write(png_bytes)
                        shape = slide.shapes.add_picture(svg_tmp, Emu(x), Emu(y), Emu(w), Emu(h))
                    finally:
                        try:
                            os.remove(svg_tmp)
                        except OSError:
                            pass
                else:
                    print(f"Warning: cannot render SVG {src} (install resvg or use programmatic shapes)")
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    tf = shape.text_frame
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = "[SVG]"
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(128, 128, 128)
            else:
                shape = slide.shapes.add_picture(src, Emu(x), Emu(y), Emu(w), Emu(h))
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

        border_spec = element.get('border')
        self._apply_border(shape, border_spec)

        # NEW: opacity for images
        opacity = element.get('opacity', 1)
        if opacity < 1:
            self._apply_opacity(shape, opacity)

        self._apply_animation(slide, shape, element)
        return shape

    def _add_table(self, slide, element):
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)
        rows_data = element.get('rows', [])
        if not rows_data:
            return None

        num_rows = len(rows_data)
        num_cols = len(rows_data[0]) if rows_data else 0

        shape = slide.shapes.add_table(num_rows, num_cols, Emu(x), Emu(y), Emu(w), Emu(h))
        table = shape.table

        style_ref = element.get('style', '')
        table_style = {}
        if isinstance(style_ref, str) and style_ref.startswith('$'):
            table_style = self.theme.table_styles.get(style_ref[1:], {})
        elif isinstance(style_ref, dict):
            table_style = style_ref

        for r_idx, row_data in enumerate(rows_data):
            for c_idx, cell_data in enumerate(row_data):
                if c_idx >= num_cols:
                    break
                cell = table.cell(r_idx, c_idx)

                if r_idx == 0:
                    fill_color = table_style.get('headerFill')
                    text_color = table_style.get('headerColor')
                    is_bold = table_style.get('headerBold', True)
                else:
                    body_fill = table_style.get('bodyFill', [])
                    fill_color = body_fill[r_idx % len(body_fill)] if body_fill else None
                    text_color = table_style.get('bodyColor')
                    is_bold = False

                if fill_color:
                    color = self.theme.resolve_color(fill_color)
                    if color:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = color

                content = cell_data.get('content', {}) if isinstance(cell_data, dict) else {'text': str(cell_data)}
                if content:
                    self._add_text_frame(cell, content)
                    font_size = table_style.get('fontSize')
                    font_family = table_style.get('fontFamily')
                    for para in cell.text_frame.paragraphs:
                        if is_bold:
                            for run in para.runs:
                                run.font.bold = True
                        if text_color:
                            color = self.theme.resolve_color(text_color)
                            if color:
                                for run in para.runs:
                                    run.font.color.rgb = color
                        if font_size:
                            for run in para.runs:
                                run.font.size = Pt(font_size)
                        if font_family:
                            for run in para.runs:
                                run.font.name = font_family.split(',')[0].strip()
                                if ',' in font_family:
                                    cjk_font = font_family.split(',')[1].strip()
                                    rPr = run._r.find(qn('a:rPr'))
                                    if rPr is None:
                                        rPr = etree.SubElement(run._r, qn('a:rPr'))
                                    ea = rPr.find(qn('a:ea'))
                                    if ea is None:
                                        ea = etree.SubElement(rPr, qn('a:ea'))
                                    ea.set('typeface', cjk_font)

        self._apply_animation(slide, shape, element)
        return shape

    def _apply_series_style(self, series_list, style):
        """Apply style dict to a list of chart series."""
        smooth = style.get('smooth')
        line_type = style.get('line')
        width = style.get('width')
        marker = style.get('marker')

        for series in series_list:
            if smooth is not None:
                series.smooth = bool(smooth)
            if width is not None:
                series.format.line.width = Pt(width)
            if line_type is not None:
                dash_map = {
                    'solid': MSO_LINE_DASH_STYLE.SOLID,
                    'dash': MSO_LINE_DASH_STYLE.DASH,
                    'dot': MSO_LINE_DASH_STYLE.ROUND_DOT,
                }
                dash = dash_map.get(line_type)
                if dash is not None:
                    series.format.line.dash_style = dash
            if marker is False:
                series.marker.style = XL_MARKER_STYLE.NONE
            elif isinstance(marker, dict) and marker.get('show') is False:
                series.marker.style = XL_MARKER_STYLE.NONE
            elif isinstance(marker, dict) and marker.get('show') is True:
                if series.marker.style == XL_MARKER_STYLE.NONE:
                    series.marker.style = XL_MARKER_STYLE.AUTOMATIC

    def _apply_chart_theme(self, chart, chartSpace, element):
        """Apply dark/light theme colors to chart via XML manipulation.
        Sets transparent backgrounds, text colors, gridline colors, and fonts."""
        ink_color = self.theme.resolve_color('$ink')
        text_color = self.theme.resolve_color('$text')
        text_light_color = self.theme.resolve_color('$textLight')
        chart_font = self.theme.get_chart_font()
        
        def _rgb_hex(color):
            return f'{color[0]:02X}{color[1]:02X}{color[2]:02X}' if color else 'FFFFFF'
        
        ink_hex = _rgb_hex(ink_color)
        text_hex = _rgb_hex(text_color)
        text_light_hex = _rgb_hex(text_light_color)
        
        # Helper to ensure a child element exists
        def _ensure_child(parent, tag):
            child = parent.find(tag)
            if child is None:
                child = etree.SubElement(parent, tag)
            return child
        
        # Helper to set text color via txPr/defRPr
        def _set_tx_color(element, color_hex, font_name=None, font_size_pt=None):
            txPr = _ensure_child(element, qn('c:txPr'))
            bodyPr = _ensure_child(txPr, qn('a:bodyPr'))
            lstStyle = _ensure_child(txPr, qn('a:lstStyle'))
            p = _ensure_child(txPr, qn('a:p'))
            pPr = _ensure_child(p, qn('a:pPr'))
            defRPr = _ensure_child(pPr, qn('a:defRPr'))
            
            # Set font color
            solidFill = _ensure_child(defRPr, qn('a:solidFill'))
            # Remove existing srgbClr/schemeClr to avoid duplicates
            for existing in list(solidFill):
                solidFill.remove(existing)
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', color_hex)
            
            if font_size_pt:
                # Size in hundredths of a point
                defRPr.set('sz', str(int(font_size_pt * 100)))
            
            if font_name:
                # Set Latin font
                latin = _ensure_child(defRPr, qn('a:latin'))
                latin.set('typeface', font_name)
                # Set East Asian font if needed
                if chart_font and ',' in chart_font:
                    cjk = chart_font.split(',')[1].strip()
                    ea = _ensure_child(defRPr, qn('a:ea'))
                    ea.set('typeface', cjk)
            
            endParaRPr = _ensure_child(p, qn('a:endParaRPr'))
            endParaRPr.set('lang', 'en-US')
        
        # Helper to set shape fill
        def _set_no_fill(element):
            spPr = _ensure_child(element, qn('c:spPr'))
            # Remove existing fill children
            for child in list(spPr):
                spPr.remove(child)
            etree.SubElement(spPr, qn('a:noFill'))
        
        # 1. ChartSpace background → transparent
        _set_no_fill(chartSpace)
        
        # 2. PlotArea background → transparent
        plotArea = chartSpace.find('.//' + qn('c:plotArea'))
        if plotArea is not None:
            _set_no_fill(plotArea)
        
        # 3. Legend text color (small size to avoid crowding)
        legend = chartSpace.find('.//' + qn('c:legend'))
        if legend is not None:
            _set_tx_color(legend, text_hex, chart_font, font_size_pt=9)
        
        # 4. Title text color
        title = chartSpace.find('.//' + qn('c:title'))
        if title is not None:
            _set_tx_color(title, ink_hex, chart_font, font_size_pt=12)
        
        # 5. Axis text colors and gridlines
        if plotArea is not None:
            for axis_tag in (qn('c:catAx'), qn('c:valAx'), qn('c:dateAx')):
                for axis in plotArea.findall(axis_tag):
                    _set_tx_color(axis, text_hex, chart_font, font_size_pt=9)
                    # Gridlines
                    for grid_tag in (qn('c:majorGridlines'), qn('c:minorGridlines')):
                        grid = axis.find(grid_tag)
                        if grid is not None:
                            g_spPr = _ensure_child(grid, qn('c:spPr'))
                            # Remove existing children
                            for child in list(g_spPr):
                                g_spPr.remove(child)
                            ln = etree.SubElement(g_spPr, qn('a:ln'))
                            ln.set('w', '6350')  # 0.5pt
                            solidFill = etree.SubElement(ln, qn('a:solidFill'))
                            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                            srgbClr.set('val', text_light_hex)
        
        # 6. Data labels text color (small size so they don't crowd the chart)
        for series in chart.series:
            dLbls = series._element.find(qn('c:dLbls'))
            if dLbls is not None:
                _set_tx_color(dLbls, ink_hex, chart_font, font_size_pt=8)
    
    def _add_chart(self, slide, element):
        """Render a real chart using python-pptx chart API.
        Supports two data formats:
          1. Simple: {labels: [...], values: [...]}
          2. Advanced: list of dicts with x_field/y_fields mapping
        """
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)
        # Support both 'type' and 'chartType' keys
        chart_type = element.get('chartType') or element.get('type', 'bar')
        data = element.get('data', [])
        x_field = element.get('x', '')
        y_fields = element.get('y', [])
        names = element.get('names', [])
        colors = element.get('colors', [])
        options = element.get('options', {})

        # Convert simple format {labels, values} to advanced format
        if isinstance(data, dict) and 'labels' in data and 'values' in data:
            labels = data['labels']
            values = data['values']
            # Build advanced data structure
            data = []
            for i, label in enumerate(labels):
                row = {'category': label}
                if i < len(values):
                    row['value'] = values[i]
                data.append(row)
            x_field = 'category'
            y_fields = ['value']
            names = ['Value']

        if not data or not x_field or not y_fields:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"[Chart: {chart_type}]"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(128, 128, 128)
            self._apply_animation(slide, shape, element)
            return shape

        if isinstance(y_fields, str):
            y_fields = [y_fields]

        direction = options.get('direction', 'vertical') if isinstance(options, dict) else 'vertical'
        if chart_type == 'bar':
            type_key = CHART_TYPE_MAP['bar'].get(direction, 'COLUMN_CLUSTERED')
        else:
            type_key = CHART_TYPE_MAP.get(chart_type, 'COLUMN_CLUSTERED')
        xl_type = getattr(XL_CHART_TYPE, type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

        from pptx.chart.data import ChartData, XyChartData, BubbleChartData

        if chart_type == 'scatter':
            chart_data = XyChartData()
            for idx, series_name in enumerate(y_fields):
                s = chart_data.add_series(names[idx] if idx < len(names) else series_name)
                for row in data:
                    x_val = row.get(x_field, 0)
                    y_val = row.get(series_name, 0)
                    s.add_data_point(x_val, y_val)
        elif chart_type == 'bubble':
            chart_data = BubbleChartData()
            size_field = element.get('size', '')
            for idx, series_name in enumerate(y_fields):
                s = chart_data.add_series(names[idx] if idx < len(names) else series_name)
                for row in data:
                    x_val = row.get(x_field, 0)
                    y_val = row.get(series_name, 0)
                    size_val = row.get(size_field, 10) if size_field else 10
                    s.add_data_point(x_val, y_val, size_val)
        else:
            chart_data = ChartData()
            categories = [str(row.get(x_field, '')) for row in data]
            chart_data.categories = categories
            for idx, series_name in enumerate(y_fields):
                display_name = names[idx] if idx < len(names) else series_name
                values = [row.get(series_name, 0) for row in data]
                chart_data.add_series(display_name, values)

        try:
            graphic_frame = slide.shapes.add_chart(
                xl_type, Emu(x), Emu(y), Emu(w), Emu(h), chart_data
            )
            chart = graphic_frame.chart

            # Apply theme colors if no explicit colors provided
            if not colors:
                palette = self.theme.get_color_palette()
                colors = [f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}" for c in palette]

            # Resolve background color for contrast checking
            bg_color = self.theme.resolve_color('$background')
            
            if colors and hasattr(chart, 'series'):
                for idx, series in enumerate(chart.series):
                    color_idx = idx % len(colors)
                    color = self.theme.resolve_color(colors[color_idx])
                    if color:
                        # Ensure color has enough contrast against background (target 3:1 for chart fills)
                        safe_color = self._ensure_contrast(color, bg_color, target=3.0)
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = safe_color
                    
                    # For single-series charts (pie, bar), color each point with a different theme color
                    if len(chart.series) == 1:
                        for p_idx, point in enumerate(series.points):
                            point_color_idx = p_idx % len(colors)
                            point_color = self.theme.resolve_color(colors[point_color_idx])
                            if point_color:
                                safe_point_color = self._ensure_contrast(point_color, bg_color, target=3.0)
                                point.format.fill.solid()
                                point.format.fill.fore_color.rgb = safe_point_color

            # Apply theme font to chart
            chart_font = self.theme.get_chart_font()
            if chart.has_title and chart.chart_title.has_text_frame:
                for para in chart.chart_title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = chart_font
            if chart.has_legend:
                try:
                    legend_tf = chart.legend.text_frame
                    for para in legend_tf.paragraphs:
                        for run in para.runs:
                            run.font.name = chart_font
                except AttributeError:
                    pass  # legend may not have text_frame

            # Apply seriesStyle settings
            series_style = element.get('seriesStyle', {})
            if series_style and hasattr(chart, 'series'):
                # Build mapping from display name to y-field name
                name_to_yfield = {}
                for idx, field_name in enumerate(y_fields):
                    display_name = names[idx] if idx < len(names) else field_name
                    name_to_yfield[display_name] = field_name

                if '*' in series_style:
                    self._apply_series_style(chart.series, series_style['*'])
                for style_key, style in series_style.items():
                    if style_key == '*':
                        continue
                    for series in chart.series:
                        yfield = name_to_yfield.get(series.name, series.name)
                        if yfield == style_key or series.name == style_key:
                            self._apply_series_style([series], style)

            title = element.get('title', '')
            if title:
                if isinstance(title, str):
                    chart.has_title = True
                    chart.chart_title.text_frame.text = title
                elif isinstance(title, dict):
                    chart.has_title = True
                    chart.chart_title.text_frame.text = title.get('text', '')

            legend = element.get('legend', True)
            if legend is False:
                chart.has_legend = False
            elif legend is True:
                chart.has_legend = True
            elif isinstance(legend, dict):
                chart.has_legend = legend.get('show', True)
                if chart.has_legend and legend.get('position'):
                    from pptx.enum.chart import XL_LEGEND_POSITION
                    pos_map = {'top': XL_LEGEND_POSITION.TOP,
                               'bottom': XL_LEGEND_POSITION.BOTTOM,
                               'left': XL_LEGEND_POSITION.LEFT,
                               'right': XL_LEGEND_POSITION.RIGHT}
                    chart.legend.position = pos_map.get(legend['position'], XL_LEGEND_POSITION.BOTTOM)

            # Auto-enable data labels by default so charts are immediately readable
            data_labels = element.get('dataLabels')
            show_labels = True
            if data_labels is False:
                show_labels = False
            elif isinstance(data_labels, dict) and data_labels.get('show') is False:
                show_labels = False
            
            if show_labels:
                for series in chart.series:
                    # Manually add data labels via XML (works for all chart types)
                    dLbls = series._element.find(qn('c:dLbls'))
                    if dLbls is None:
                        dLbls = etree.SubElement(series._element, qn('c:dLbls'))
                    
                    # Clear existing children
                    for child in list(dLbls):
                        dLbls.remove(child)
                    
                    if chart_type == 'pie':
                        etree.SubElement(dLbls, qn('c:showVal')).set('val', '0')
                        etree.SubElement(dLbls, qn('c:showPercent')).set('val', '1')
                    else:
                        etree.SubElement(dLbls, qn('c:showVal')).set('val', '1')
                        etree.SubElement(dLbls, qn('c:showPercent')).set('val', '0')
                    
                    etree.SubElement(dLbls, qn('c:showCatName')).set('val', '0')
                    etree.SubElement(dLbls, qn('c:showSerName')).set('val', '0')
                    etree.SubElement(dLbls, qn('c:showLegendKey')).set('val', '0')
            
            # Apply theme styling (background, text colors, gridlines)
            self._apply_chart_theme(chart, chart._chartSpace, element)

            if chart_type == 'pie' and isinstance(options, dict):
                inner_radius = options.get('innerRadius')
                if inner_radius and hasattr(chart, 'plot_type'):
                    pass

            return graphic_frame
        except (ValueError, KeyError, TypeError, AttributeError, RuntimeError) as e:
            print(f"Error rendering chart: {e}")
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"[Chart: {chart_type}]"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(128, 128, 128)
            self._apply_animation(slide, shape, element)
            return shape

    def _add_icon(self, slide, element):
        """Render icon as a freeform shape from path library, or fallback to placeholder text."""
        bounds = element.get('bounds', [0, 0, 100, 100])
        x, y, w, h = self._to_emu(*bounds)
        icon_name = element.get('iconName', 'fas:question')

        # Try to resolve icon path from library
        try:
            from icon_paths import get_icon_path
            path_str = get_icon_path(icon_name)
        except (ImportError, NameError, AttributeError):
            path_str = None

        if path_str:
            # Render as freeform shape
            icon_element = dict(element)
            icon_element['path'] = path_str
            icon_element['shapeName'] = 'custom'
            result = self._add_freeform_shape(slide, icon_element)
            if result is not None:
                return result

        # Fallback: render a simple geometric placeholder shape with no text
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(x), Emu(y), Emu(w), Emu(h)
        )
        fill_spec = element.get('fill')
        if fill_spec:
            self._apply_fill(shape, fill_spec)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme.resolve_color('$primary') or RGBColor(128, 128, 128)
        shape.line.fill.background()

        self._apply_animation(slide, shape, element)
        return shape

    def _set_background(self, slide, bg_spec):
        if not bg_spec:
            return
        bg_type = bg_spec.get('type')
        if bg_type == 'solid':
            color = self.theme.resolve_color(bg_spec.get('color'))
            if color:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = color
        elif bg_type == 'gradient':
            stops = bg_spec.get('stops', [])
            if stops:
                background = slide.background
                fill = background.fill
                fill.gradient()
                # For slide background, python-pptx gradient support is limited
                # Use first stop as fore_color
                color = self.theme.resolve_color(stops[0].get('color'))
                if color:
                    fill.fore_color.rgb = color
                if len(stops) > 1:
                    color2 = self.theme.resolve_color(stops[-1].get('color'))
                    if color2:
                        fill.back_color.rgb = color2
        elif bg_type == 'image':
            src = bg_spec.get('src', '')
            if src:
                if not os.path.isabs(src):
                    src = os.path.join(self.base_dir, src)
                if os.path.exists(src):
                    # Add as full-slide picture behind everything
                    x, y, w, h = self._to_emu(0, 0, self.pptd_w, self.pptd_h)
                    pic_shape = slide.shapes.add_picture(src, Emu(x), Emu(y), Emu(w), Emu(h))

                    # Handle opacity on background image
                    opacity = bg_spec.get('opacity', 1)
                    if opacity < 1:
                        self._apply_opacity(pic_shape, opacity)

                    # Handle mask overlay
                    mask = bg_spec.get('mask')
                    if mask:
                        mask_color = self.theme.resolve_color(mask.get('color', '#000000'))
                        mask_opacity = mask.get('opacity', 0.5)
                        mask_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
                        )
                        mask_shape.fill.solid()
                        mask_shape.fill.fore_color.rgb = mask_color or RGBColor(0, 0, 0)
                        self._apply_opacity(mask_shape, mask_opacity)
                        mask_shape.line.fill.background()

    def _infer_layer(self, element):
        """Infer z-order layer for an element when not explicitly set.
        
        Layer conventions:
          -2 = background patterns (full-screen, low opacity textures)
          -1 = background color blocks (full-screen solid fills)
           0 = mid-ground containers, cards (default)
           1 = content (text, charts, tables, images)
           2 = foreground decorations (accent lines, corner brackets, labels)
        """
        if 'layer' in element:
            return element['layer']
        
        etype = element.get('elementType', '')
        bounds = element.get('bounds', [0, 0, 0, 0])
        opacity = element.get('opacity', 1)
        
        x, y, w, h = bounds
        is_near_fullscreen = w >= 1200 and h >= 680
        is_small = w <= 100 and h <= 20

        # Decorative text: large numerals at low opacity
        if etype == 'text' and opacity < 0.3:
            font_size = element.get('content', {}).get('fontSize', 0)
            if font_size > 100:
                return -1
            if is_near_fullscreen:
                return -2

        # Background patterns: full-screen + low opacity
        if etype == 'shape' and is_near_fullscreen and opacity < 0.3:
            return -2

        # Background color blocks: full-screen solid fills
        if etype == 'shape' and is_near_fullscreen:
            fill = element.get('fill', {})
            if fill.get('type') in ('solid', 'gradient'):
                return -1

        # Foreground decorations: tiny shapes (accent lines, small badges)
        if etype == 'shape' and is_small:
            return 2

        # Content elements
        if etype in ('text', 'chart', 'table', 'image'):
            return 1
        
        # Default: mid-ground containers
        return 0

    def convert_page(self, slide, page_data):
        """Convert a single .page data dict to a slide."""
        # Support both 'bgColor' (string, e.g. "$background" or "#FFFFFF")
        # and 'background' (dict, e.g. {type: solid, color: "..."})
        bg_color = page_data.get('bgColor')
        if bg_color:
            resolved = self.theme.resolve_color(bg_color)
            if resolved:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = resolved
        
        bg = page_data.get('background')
        if bg:
            self._set_background(slide, bg)

        # NEW: Handle page-level decoration directives
        decorations = page_data.get('decorations', [])
        for deco in decorations:
            dtype = deco.get('type')
            if dtype == 'corner_brackets':
                self._draw_corner_brackets(
                    slide, deco.get('bounds', [0, 0, 1280, 720]),
                    deco.get('color', '#000000'),
                    deco.get('strokeWidth', 4),
                    deco.get('size', 24),
                    deco.get('style', 'L')
                )
            elif dtype == 'hatch':
                self._draw_hatch_pattern(
                    slide, deco.get('bounds', [0, 0, 1280, 720]),
                    deco.get('color', '#000000'),
                    deco.get('angle', 45),
                    deco.get('spacing', 20),
                    deco.get('lineWidth', 2),
                    deco.get('opacity', 0.06)
                )
            elif dtype == 'scanlines':
                self._draw_scanlines(
                    slide, deco.get('bounds', [0, 0, 1280, 720]),
                    deco.get('color', '#000000'),
                    deco.get('spacing', 4),
                    deco.get('lineHeight', 2),
                    deco.get('opacity', 0.04)
                )
            elif dtype == 'grid':
                self._draw_grid(
                    slide, deco.get('bounds', [0, 0, 1280, 720]),
                    deco.get('color', '#000000'),
                    deco.get('gridSize', 40),
                    deco.get('lineWidth', 1),
                    deco.get('opacity', 0.07)
                )

        # NEW: Sort elements by inferred layer to prevent occlusion
        elements = page_data.get('elements', [])
        sorted_elements = sorted(elements, key=lambda e: self._infer_layer(e))
        
        # Report any reordering for debugging
        if len(elements) > 1:
            original_order = [e.get('elementId', '?') for e in elements]
            sorted_order = [e.get('elementId', '?') for e in sorted_elements]
            if original_order != sorted_order:
                reordered = [(o, s) for o, s in zip(original_order, sorted_order) if o != s]
                if reordered:
                    print(f"  Layer reorder: {len(reordered)} elements reordered")

        for element in sorted_elements:
            etype = element.get('elementType')
            try:
                if etype == 'shape':
                    self._add_shape(slide, element)
                elif etype == 'text':
                    self._add_text(slide, element)
                elif etype == 'image':
                    self._add_image(slide, element)
                elif etype == 'table':
                    self._add_table(slide, element)
                elif etype == 'chart':
                    self._add_chart(slide, element)
                elif etype == 'icon':
                    self._add_icon(slide, element)
            except (ValueError, KeyError, TypeError, AttributeError, RuntimeError) as e:
                eid = element.get('elementId', 'unknown')
                print(f"Error processing element {eid} ({etype}): {e}")

    def _draw_grain_texture(self, slide, bounds, color, density=0.03, dot_size=2):
        """Draw grain/noise texture using many small dots.
        Uses a deterministic pseudo-random pattern so output is reproducible."""
        bx, by, bw, bh = bounds
        # Use a simple hash-based pseudo-random for reproducibility
        def prng(seed):
            while True:
                seed = (seed * 1103515245 + 12345) & 0x7fffffff
                yield seed / 0x7fffffff

        rng = prng(42)
        # Calculate number of dots based on area and density
        area = bw * bh
        num_dots = max(20, int(area * density))

        for _ in range(num_dots):
            rx = bx + next(rng) * bw
            ry = by + next(rng) * bh
            rw = dot_size + next(rng) * dot_size
            rh = dot_size + next(rng) * dot_size
            elem = {
                'elementType': 'shape',
                'shapeName': 'rect',
                'bounds': [rx, ry, rw, rh],
                'fill': {'type': 'solid', 'color': color},
                'opacity': 0.15 + next(rng) * 0.25,
            }
            self._add_shape(slide, elem)

    def _add_page_number(self, slide, page_num, total_pages, config):
        """Add automatic page number to slide."""
        position = config.get('position', 'bottom-right')
        color = config.get('color', '#000000')
        font_size = config.get('fontSize', 10)
        font_family = config.get('fontFamily', 'Space Grotesk')
        format_str = config.get('format', '{n}')
        offset_x = config.get('offsetX', 40)
        offset_y = config.get('offsetY', 30)

        text = format_str.format(n=page_num, total=total_pages)

        # Calculate position
        if position == 'bottom-right':
            x = 1280 - offset_x - 60
            y = 720 - offset_y - 20
        elif position == 'bottom-left':
            x = offset_x
            y = 720 - offset_y - 20
        elif position == 'bottom-center':
            x = 640 - 30
            y = 720 - offset_y - 20
        elif position == 'top-right':
            x = 1280 - offset_x - 60
            y = offset_y
        elif position == 'top-left':
            x = offset_x
            y = offset_y
        else:
            x = 1280 - offset_x - 60
            y = 720 - offset_y - 20

        elem = {
            'elementType': 'text',
            'bounds': [x, y, 60, 20],
            'content': {
                'text': text,
                'fontSize': font_size,
                'color': color,
                'fontFamily': font_family,
            },
            'wrap': False,
        }
        self._add_text(slide, elem)

    def convert(self, output_path=None):
        if output_path is None:
            base, _ = os.path.splitext(self.pptd_path)
            output_path = base + '.pptx'

        pages = self.pptd.get('pages', [])
        blank_layout = self.prs.slide_layouts[6]

        # Global page number config
        page_number_config = self.pptd.get('pageNumber')
        total_pages = len(pages)

        for page_idx, page_ref in enumerate(pages):
            page_path = os.path.join(self.base_dir, page_ref)
            if not os.path.exists(page_path):
                print(f"Warning: page file not found: {page_path}")
                continue

            with open(page_path, encoding='utf-8') as f:
                page_data = yaml.safe_load(f)

            slide = self.prs.slides.add_slide(blank_layout)
            self.convert_page(slide, page_data)

            # Add automatic page number if configured
            if page_number_config:
                page_num = page_idx + 1
                self._add_page_number(slide, page_num, total_pages, page_number_config)

        self.prs.save(output_path)
        print(f"Saved: {output_path}")
        return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTD presentations to PPTX files",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptd2pptx.py presentation.pptd
  python pptd2pptx.py presentation.pptd output.pptx
  python pptd2pptx.py presentation.pptd --verbose
        """,
    )
    parser.add_argument("pptd", help="Path to the .pptd master file")
    parser.add_argument("output", nargs="?", help="Output .pptx path (default: same base name)")
    parser.add_argument("-v", "--verbose", action="store_true", help="Print detailed conversion info")
    parser.add_argument("--check", action="store_true", help="Run validation before conversion")

    args = parser.parse_args()

    pptd_path = args.pptd
    output_path = args.output

    if not os.path.exists(pptd_path):
        print(f"Error: File not found: {pptd_path}", file=sys.stderr)
        sys.exit(1)

    if args.check:
        try:
            from check_pptd import Checker
            checker = Checker(pptd_path)
            result = checker.run()
            if result != 0:
                print("Validation failed. Use --verbose for details.", file=sys.stderr)
                sys.exit(1)
            print("Validation passed.")
        except (ImportError, OSError, RuntimeError) as e:
            print(f"Warning: Could not run validation: {e}", file=sys.stderr)

    converter = PPTDConverter(pptd_path)
    if args.verbose:
        print(f"Converting: {pptd_path}")
        print(f"  Pages: {len(converter.pptd.get('pages', []))}")
        print(f"  Theme colors: {list(converter.pptd.get('theme', {}).get('colors', {}).keys())}")

    try:
        result = converter.convert(output_path)
        print(f"Saved: {result}")
    except Exception as e:
        print(f"Error: Conversion failed: {e}", file=sys.stderr)
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
